from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

class PowerPoint:
    def SubElement(self, parent, tagname, **kwargs): # Para personalizar as características do slide
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _set_shape_transparency(self, shape, alpha): # Para personalizar as características de transparência do slide
        ts = shape.fill._xPr.solidFill
        sF = ts.get_or_change_to_srgbClr()
        self.SubElement(sF, 'a:alpha', val=str(alpha))

    def generatePptx(self, title, author, lyrics):
        left = top = 0
        width = Inches(16)
        height = Inches(9)

        pptx = Presentation()
        pptx.slide_width = width
        pptx.slide_height = height

        fileName = title + " - " + author + ".pptx"

        #
        ## Criar um slide incial para o título e autor/banda, e definir seu tipo de slide para 6 (vazio)
        layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(layout)

        #
        ## Adicionar forma preta no fundo
        shapes = slide.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )

        #
        ## Adicionar transparência à forma
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        self._set_shape_transparency(shape, 50000)

        line = shape.line
        line.fill.background()

        #
        ## Adicionar caixa de texto
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

        #
        ## Escrever o título da música e definir as características do seu texto
        p = tf.paragraphs[0]
        p.text = title
        # self._set_p_border(p)
        p.font.bold = True
        p.font.size = Pt(80)
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255,255,255)

        #
        ## Escrever o autor/banda da música e definir as características do seu texto
        p = tf.add_paragraph()
        p.text = author
        p.font.bold = True
        p.font.italic = True
        p.font.size = Pt(50)
        p.alignment = PP_ALIGN.CENTER
        p.word_wrap = True
        p.font.color.rgb = RGBColor(255,255,255)

        for i in range(len(lyrics)): # Adicionar um slide novo para cada parágrafo da letra da música
            #
            ## Criação do slide e definição do seu tipo
            layout = pptx.slide_layouts[6]
            slide = pptx.slides.add_slide(layout)

            #
            ## Criação de uma forma para o fundo
            shapes = slide.shapes
            shape = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )

            #
            ## Adicionar transparência à forma
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            self._set_shape_transparency(shape, 50000)

            line = shape.line
            line.fill.background()

            #
            ## Adicionar caixa de texto
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True

            #
            ## Escrever o parágrafo no slide e definir a característica do texto
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.word_wrap = True
            p.text = lyrics[i].replace('\n\n', '')
            p.font.bold = True
            p.font.size = Pt(60)
            p.font.color.rgb = RGBColor(255,255,255)
        pptx.save(fileName)

        return fileName
