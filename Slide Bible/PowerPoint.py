from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

class PowerPoint:
    def generatePptx(self, reference):
        fileName = "biblia.pptx"
        book_chapter = reference[0]

        left = top = 0
        width = Inches(16)
        height = Inches(9)

        pptx = Presentation()
        pptx.slide_width = width
        pptx.slide_height = height


        for i in range(len(reference)): # Adicionar um slide novo para cada versículo
            if (i == 0): continue # Pular primeiro elemento
            verse = reference[i]

            #
            ## Criação do slide e definição do seu tipo, fundo preto
            layout = pptx.slide_layouts[6]
            slide = pptx.slides.add_slide(layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            #
            ## Adicionar caixa de texto para referência
            txBox = slide.shapes.add_textbox(Inches(0.1), Inches(8.9), width, Inches(0.5)) if ":" in book_chapter else \
                slide.shapes.add_textbox(Inches(0.1), Inches(8.4), width, Inches(0.5)) # left, top, width, height
            text_frame = txBox.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

            #
            ## Escrever o parágrafo no slide e definir a característica do texto de referência
            p = text_frame.paragraphs[0]
            p.text = book_chapter
            p.font.bold = True
            p.font.size = Pt(30)
            p.font.color.rgb = RGBColor(255, 255, 255)

            #
            ## Adicionar caixa de texto para versículo
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(15), Inches(7.7)) # left, top, width, height
            text_frame = txBox.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.word_wrap = True

            #
            ## Escrever o parágrafo no slide e definir a característica do texto de versículo
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.JUSTIFY_LOW
            p.word_wrap = True
            p.text = verse
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        pptx.save(fileName)

        return fileName
